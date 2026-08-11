import os
import re
import json
import logging
from typing import Any, List, Dict, Tuple
from pathlib import Path
from urllib.parse import urlparse

import pandas as pd
from pypdf import PdfReader
import docx2txt
from bs4 import BeautifulSoup

from ..html_text_utils import enrich_contact_links
from ..pdf_text_cleaner import normalize_pdf_extracted_text

logger = logging.getLogger(__name__)

# --- Chunk config ---
CHUNK_OVERLAP = 200
CHUNK_SIZES = {
    ".html": 1400,
    ".htm": 1400,
    ".pdf": 1200,
    ".docx": 1200,
    ".doc": 1200,
    ".pptx": 1000,
    ".txt": 1500,
    ".md": 1500,
}

# Mis-suffixed OOXML/ZIP bytes read as text produce these markers.
_OOXML_GARBAGE_MARKERS = (
    "[Content_Types].xml",
    "ppt/slides/",
    "ppt/presentation.xml",
    "word/document.xml",
    "xl/workbook.xml",
    "_rels/.rels",
)

# Lone Unicode surrogates (U+D800..U+DFFF) leak out of pypdf when a PDF has
# malformed text streams. Chroma/JSON serialization rejects them with
# "surrogates not allowed", which crashes the entire ingest and leaves the
# document permanently un-indexed. Strip them before anything else touches
# the extracted text.
_SURROGATE_RE = re.compile(r"[\ud800-\udfff]")


def _strip_surrogates(text: str) -> str:
    if not text:
        return ""
    return _SURROGATE_RE.sub("", text)


# Lazily-built sentence splitter so the import cost is paid once.
_sentence_splitter_cache: Dict[Tuple[int, int], Any] = {}


def _get_sentence_splitter(chunk_size: int, overlap: int):
    """Build a sentence-aware splitter whose chunk_size/overlap are in characters.

    llama-index's SentenceSplitter measures chunk_size in *tokens* by default
    (via tiktoken). Our existing CHUNK_SIZES table is in *characters* (1200
    for PDF, etc.), so we plug in ``tokenizer=list`` to make each character
    count as one token. This preserves the historical chunk-granularity
    while gaining sentence/paragraph awareness on top.
    """
    key = (int(chunk_size), int(overlap))
    splitter = _sentence_splitter_cache.get(key)
    if splitter is not None:
        return splitter
    try:
        from llama_index.core.node_parser import SentenceSplitter

        splitter = SentenceSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            tokenizer=list,
        )
    except Exception as exc:
        logger.warning(
            "SentenceSplitter unavailable (%s); falling back to char-based chunking", exc
        )
        splitter = None
    _sentence_splitter_cache[key] = splitter
    return splitter


# ---------------- Helpers ----------------
def normalize_url(url: str) -> str:
    """Normalize URL by removing extra whitespace, backticks, and fixing double slashes."""
    if not url:
        return ""

    url = url.strip().strip("`").strip()
    if not url:
        return ""

    # Skip non-HTTP URLs
    if not url.startswith(("http://", "https://", "file://")):
        return url

    try:
        parsed = urlparse(url)
        path = re.sub(r"/+", "/", parsed.path) if parsed.path else ""

        scheme = parsed.scheme
        netloc = parsed.netloc
        params = parsed.params
        query = parsed.query
        fragment = parsed.fragment

        normalized = f"{scheme}://{netloc}{path}"
        if params:
            normalized += f";{params}"
        if query:
            normalized += f"?{query}"
        if fragment:
            normalized += f"#{fragment}"

        return normalized
    except Exception as e:
        logger.warning(f"URL normalization failed for '{url}': {e}")
        return url


def chunk_text(text: str, chunk_size: int, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into chunks aligned to sentence boundaries when possible.

    Uses llama-index's SentenceSplitter (already a project dep) to avoid
    mid-sentence / mid-word cuts that produce noisy embeddings. Falls back to
    the previous character-window splitter if llama-index is unavailable.
    """
    safe_text = _strip_surrogates(text or "")
    if not safe_text.strip():
        return []

    splitter = _get_sentence_splitter(chunk_size, overlap)
    if splitter is not None:
        try:
            parts = splitter.split_text(safe_text)
            return [p.strip() for p in parts if p and p.strip()]
        except Exception as exc:
            logger.warning("SentenceSplitter failed (%s); falling back to char chunking", exc)

    chunks: List[str] = []
    start = 0
    text_len = len(safe_text)
    while start < text_len:
        end = start + chunk_size
        chunks.append(safe_text[start:end].strip())
        if end >= text_len:
            break
        start = end - overlap
    return [c for c in chunks if c]


def looks_like_ooxml_garbage(text: str) -> bool:
    """True when text looks like raw OOXML/ZIP bytes decoded as UTF-8."""
    sample = (text or "")[:8000]
    if not sample.strip():
        return False
    head = sample[:64]
    if head.startswith("PK") or "\x00" in sample[:200]:
        marker_hits = sum(1 for m in _OOXML_GARBAGE_MARKERS if m in sample)
        if marker_hits >= 1:
            return True
    marker_hits = sum(1 for m in _OOXML_GARBAGE_MARKERS if m in sample)
    return marker_hits >= 2


def _extract_pptx_slide_texts(filepath: str) -> List[Tuple[int, str]]:
    """Return [(slide_number, text), ...] from a .pptx file."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(filepath)
    slides: List[Tuple[int, str]] = []

    def _shape_text(shape) -> List[str]:
        parts: List[str] = []
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = (shape.text_frame.text or "").strip()
            if text:
                parts.append(text)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE and getattr(shape, "has_table", False):
            try:
                table = shape.table
                for row in table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        parts.append(row_text)
            except Exception:
                pass
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    parts.extend(_shape_text(child))
            except Exception:
                pass
        return parts

    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            try:
                parts.extend(_shape_text(shape))
            except Exception:
                continue
        if slide.has_notes_slide:
            try:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(notes)
            except Exception:
                pass
        combined = "\n".join(p for p in parts if p and p.strip()).strip()
        if combined:
            slides.append((slide_idx, _strip_surrogates(combined)))
    return slides


def chunk_size_for_crawled_url(url: str) -> int:
    """Chunk size (chars) for a crawled URL — aligned with file ingest in this module."""
    path = urlparse(url or "").path.lower().split("?")[0]
    if path.endswith(".pdf"):
        return CHUNK_SIZES[".pdf"]
    if path.endswith((".html", ".htm")):
        return CHUNK_SIZES[".html"]
    if path.endswith((".docx", ".doc")):
        return CHUNK_SIZES[".docx"]
    if path.endswith(".pptx"):
        return CHUNK_SIZES[".pptx"]
    if path.endswith((".md",)):
        return CHUNK_SIZES[".md"]
    if path.endswith((".txt",)):
        return CHUNK_SIZES[".txt"]
    return CHUNK_SIZES.get(".txt", 900)


def chunks_for_crawled_document(url: str, text: str) -> List[str]:
    """Split crawled page text into embedding-sized chunks (same rules as uploaded files)."""
    safe_text = text or ""
    if urlparse(url or "").path.lower().split("?")[0].endswith(".pdf"):
        safe_text = normalize_pdf_extracted_text(safe_text)
    return chunk_text(safe_text, chunk_size_for_crawled_url(url))


def _as_primitive(x: Any) -> Any:
    """Serialize metadata values to safe primitives (string/number/bool)."""
    if x is None:
        return ""
    if isinstance(x, (str, int, float, bool)):
        return x
    try:
        return json.dumps(x)
    except Exception:
        return str(x)


def _extract_title_from_row(row: pd.Series) -> str:
    for col in ["title", "Title", "name", "Name"]:
        if col in row and pd.notna(row[col]):
            return str(row[col]).strip()
    return ""


def _detect_url_column(df: pd.DataFrame) -> List[str]:
    """Return candidate url column names."""
    candidates = [
        "url",
        "link",
        "href",
        "website",
        "source",
        "reference",
        "page_url",
        "pageUrl",
        "canonical_url",
        "canonicalUrl",
        "permalink",
        
    ]
    cols = []
    for col in df.columns:
        col_lower = col.lower().strip()
        if col_lower in candidates or any(p in col_lower for p in candidates):
            cols.append(col)
    return cols


def _extract_url_from_row(
    row: pd.Series, df: pd.DataFrame, filepath: str, row_index: int
) -> str:
    url_columns = _detect_url_column(df)
    for col in url_columns:
        if pd.notna(row[col]):
            url_val = str(row[col]).strip()
            if url_val:
                url = normalize_url(url_val)
                if not url.startswith(("http://", "https://", "file://")):
                    url = f"file://{os.path.basename(filepath)}#row-{row_index + 1}"
                return url
    return f"file://{os.path.basename(filepath)}#row-{row_index + 1}"


def _extract_keywords_from_row(row: pd.Series) -> str:
    for col in ["keywords", "keyword", "tags", "Keywords"]:
        if col in row and pd.notna(row[col]):
            return str(row[col]).strip()
    return ""


def _extract_content_from_row(row: pd.Series) -> str:
    for col in [
        "meaningful_content",
        "Meaningful_Content",
        "meaningfulContent",
        "content",
        "text",
        "description",
        "text_content",
    ]:
        if col in row and pd.notna(row[col]):
            return str(row[col]).strip()
    # fallback: concatenation
    return " ".join([str(v) for v in row.values if pd.notna(v)])


def extract_text_from_file(filepath: str) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Return (text_chunks, metadata_per_chunk) for a given file."""
    ext = os.path.splitext(filepath)[1].lower()
    texts: List[str] = []
    metadata_list: List[Dict[str, Any]] = []

    if ext == ".pdf":
        pdf_reader = PdfReader(filepath)
        raw_text = "\n".join([page.extract_text() or "" for page in pdf_reader.pages])
        raw_text = _strip_surrogates(raw_text)
        raw_text = normalize_pdf_extracted_text(raw_text)
        texts = chunk_text(raw_text, CHUNK_SIZES[".pdf"])

    elif ext in [".docx", ".doc"]:
        raw_text = _strip_surrogates(docx2txt.process(filepath) or "")
        if looks_like_ooxml_garbage(raw_text):
            raise ValueError(
                "DOCX extraction produced OOXML/ZIP garbage — file may be mislabeled"
            )
        texts = chunk_text(raw_text, CHUNK_SIZES[".docx"])

    elif ext == ".pptx":
        slide_texts = _extract_pptx_slide_texts(filepath)
        if not slide_texts:
            raise ValueError("No text extracted from PPTX")
        combined_preview = "\n".join(t for _, t in slide_texts[:5])
        if looks_like_ooxml_garbage(combined_preview):
            raise ValueError(
                "PPTX extraction produced OOXML/ZIP garbage — file may be mislabeled"
            )
        title = os.path.splitext(os.path.basename(filepath))[0]
        texts = []
        metadata_list = []
        for slide_num, slide_text in slide_texts:
            for local_idx, chunk in enumerate(
                chunk_text(slide_text, CHUNK_SIZES[".pptx"])
            ):
                texts.append(chunk)
                metadata_list.append(
                    {
                        "title": _as_primitive(title),
                        "url": _as_primitive(f"file://{os.path.basename(filepath)}"),
                        "keywords": "",
                        "chunk_index": len(texts) - 1,
                        "slide_number": slide_num,
                        "page": slide_num,
                        "slide_chunk_index": local_idx,
                    }
                )

    elif ext == ".ppt":
        raise ValueError(
            "Unsupported file type: .ppt (legacy PowerPoint). Convert to .pptx and re-upload."
        )

    elif ext in [".txt", ".md"]:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            raw_text = _strip_surrogates(f.read())
        if looks_like_ooxml_garbage(raw_text):
            raise ValueError(
                "Text extraction looks like OOXML/ZIP bytes — refusing to index binary garbage"
            )
        texts = chunk_text(raw_text, CHUNK_SIZES.get(ext, 900))

    elif ext in [".html", ".htm"]:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
            enrich_contact_links(soup)
            raw_text = _strip_surrogates(soup.get_text(separator="\n"))

            urls_from_html: List[str] = []
            for link in soup.find_all("a", href=True):
                href = link.get("href")
                if not href:
                    continue
                normalized_href = normalize_url(href)
                if normalized_href.startswith(("http://", "https://", "file://")):
                    urls_from_html.append(normalized_href)

            texts = chunk_text(raw_text, CHUNK_SIZES[".html"])
            if urls_from_html:
                urls_text = "URLs found in document: " + ", ".join(sorted(set(urls_from_html)))
                texts.append(urls_text)

    elif ext == ".csv":
        df = pd.read_csv(filepath, on_bad_lines="skip", dtype=str)
        texts = []
        metadata_list = []

        for idx, row in df.iterrows():
            title = _extract_title_from_row(row)
            url = _extract_url_from_row(row, df, filepath, idx)
            keywords = _extract_keywords_from_row(row)
            content = _extract_content_from_row(row)

            content = _strip_surrogates(content)
            if not content.strip():
                continue

            content_chunks = chunk_text(content, CHUNK_SIZES.get(".txt", 900))
            for chunk_idx, chunk_text_content in enumerate(content_chunks):
                texts.append(chunk_text_content)
                metadata_list.append(
                    {
                        "title": _as_primitive(title),
                        "url": _as_primitive(url),
                        "keywords": _as_primitive(keywords),
                        "chunk_index": chunk_idx,
                        "row_index": idx,
                    }
                )
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    logger.info(f"Extracted {len(texts)} chunks from {filepath}")

    # For non-CSV / non-PPTX: generate basic metadata per chunk
    if ext not in (".csv", ".pptx"):
        metadata_list = []
        title = os.path.splitext(os.path.basename(filepath))[0]
        url_pattern = r"https?://[^\s,<>'\")]+|www\.[^\s,<>'\")]+"

        for chunk_idx, text in enumerate(texts):
            found_urls = re.findall(url_pattern, text)
            chunk_url = ""
            if found_urls:
                chunk_url = normalize_url(found_urls[0])
                if not chunk_url.startswith(("http://", "https://", "file://")):
                    chunk_url = ""
            if not chunk_url:
                chunk_url = f"file://{os.path.basename(filepath)}"

            metadata_list.append(
                {
                    "title": _as_primitive(title),
                    "url": _as_primitive(chunk_url),
                    "keywords": "",
                    "chunk_index": chunk_idx,
                }
            )

    return texts, metadata_list


def batch_iterate(lst: List[Any], size: int):
    for i in range(0, len(lst), size):
        yield lst[i : i + size]
