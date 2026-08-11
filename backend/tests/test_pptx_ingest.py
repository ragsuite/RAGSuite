"""PPTX extraction, Drive MIME, reindex suffix, and OOXML garbage guards."""
from __future__ import annotations

import io
import zipfile
from pathlib import Path
from unittest.mock import MagicMock

import pytest

from app.services.connectors.google_drive import (
    _download_file_content,
    _normalize_doc_type,
)
from app.services.rag.utils_rag import (
    extract_text_from_file,
    looks_like_ooxml_garbage,
)
from app.services.reindex_service import reindex_temp_suffix_for_uploaded_doc


def _make_minimal_pptx(path: Path, slide_text: str = "Online Voting System overview") -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_text
    body = slide.placeholders[1]
    body.text = "Secure authentication and encrypted ballots for digital elections."
    prs.save(str(path))


def test_looks_like_ooxml_garbage_detects_zip_text():
    assert looks_like_ooxml_garbage("PK\x03\x04 [Content_Types].xml ppt/slides/slide1.xml")
    assert not looks_like_ooxml_garbage("Online Voting System uses encrypted ballots.")


def test_extract_pptx_yields_readable_slide_text(tmp_path: Path):
    pptx_path = tmp_path / "voting.pptx"
    _make_minimal_pptx(pptx_path, "Online Voting System")
    texts, metas = extract_text_from_file(str(pptx_path))
    joined = "\n".join(texts)
    assert "Online Voting System" in joined
    assert "encrypted ballots" in joined.lower() or "Secure authentication" in joined
    assert "[Content_Types].xml" not in joined
    assert all(m.get("slide_number") == 1 for m in metas)


def test_extract_txt_refuses_ooxml_garbage(tmp_path: Path):
    junk = tmp_path / "fake.txt"
    junk.write_bytes(b"PK\x03\x04\n[Content_Types].xml\nppt/presentation.xml\n")
    with pytest.raises(ValueError, match="OOXML/ZIP"):
        extract_text_from_file(str(junk))


def test_extract_ppt_legacy_rejected(tmp_path: Path):
    legacy = tmp_path / "old.ppt"
    legacy.write_bytes(b"\xd0\xcf\x11\xe0")  # OLE header stub
    with pytest.raises(ValueError, match=r"\.ppt"):
        extract_text_from_file(str(legacy))


def test_normalize_drive_pptx_type_is_not_bin():
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert len(mime) > 50
    shortened = _normalize_doc_type(mime, ".pptx")
    assert shortened == "application/PPTX"
    assert "BIN" not in shortened


def test_download_file_content_maps_pptx_ext(monkeypatch):
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    file_meta = {"id": "abc", "mimeType": mime, "size": 10, "name": "deck.pptx"}

    class _FakeRequest:
        pass

    class _FakeFiles:
        def get_media(self, fileId):
            assert fileId == "abc"
            return _FakeRequest()

    class _FakeService:
        def files(self):
            return _FakeFiles()

    class _FakeDownloader:
        def __init__(self, buffer, request):
            self.buffer = buffer
            self._done = False

        def next_chunk(self):
            if not self._done:
                self.buffer.write(b"PK-fake")
                self._done = True
                return {}, False
            return {}, True

    monkeypatch.setattr(
        "googleapiclient.http.MediaIoBaseDownload",
        _FakeDownloader,
    )
    content, ext = _download_file_content(_FakeService(), file_meta, max_size_bytes=1024)
    assert ext == ".pptx"
    assert content == b"PK-fake"


def test_reindex_suffix_pptx_from_title_and_bin_type():
    doc = MagicMock()
    doc.title = "ONLINE VOTING SYSTEM.pptx"
    doc.type = "application/BIN"
    # Minimal ZIP containing ppt/ marker
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("[Content_Types].xml", "<Types/>")
        zf.writestr("ppt/presentation.xml", "<p:presentation/>")
    raw = buf.getvalue()
    assert reindex_temp_suffix_for_uploaded_doc(doc, raw) == ".pptx"


def test_reindex_suffix_refuses_unknown_zip_as_txt():
    doc = MagicMock()
    doc.title = "mystery.bin"
    doc.type = "application/BIN"
    raw = b"PK\x03\x04" + b"\x00" * 100
    with pytest.raises(ValueError, match="refusing .txt fallback"):
        reindex_temp_suffix_for_uploaded_doc(doc, raw)
